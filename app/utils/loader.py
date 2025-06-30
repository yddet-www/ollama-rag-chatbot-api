import os
from bs4 import BeautifulSoup
import docx
import PyPDF2
from pptx import Presentation
from langchain_core.documents import Document

# Load and parse supported document formats with filename as context
def load_documents(directory: str):
    documents = []
    for filename in os.listdir(directory):
        filepath = os.path.join(directory, filename)
        try:
            file_text = ""

            if filename.endswith(".txt"):
                with open(filepath, "r", encoding="utf-8") as f:
                    file_text = f.read()

            elif filename.endswith(".xml"):
                with open(filepath, "r", encoding="utf-8") as f:
                    content = f.read()
                soup = BeautifulSoup(content, "lxml-xml")
                file_text = soup.get_text(separator=" ", strip=True)

            elif filename.endswith(".docx"):
                doc = docx.Document(filepath)
                file_text = "\n".join([p.text for p in doc.paragraphs])

            elif filename.endswith(".pdf"):
                with open(filepath, "rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    file_text = "\n".join(
                        page.extract_text() for page in reader.pages if page.extract_text()
                    )

            elif filename.endswith(".pptx"):
                prs = Presentation(filepath)
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            texts.append(shape.text)
                file_text = "\n".join(texts)

            elif filename.endswith(".csv"):
                with open(filepath, "r", encoding="utf-8") as f:
                    file_text = f.read()

            if file_text.strip():  # Avoid empty documents
                filename_only = os.path.basename(filepath)
                filename_base = os.path.splitext(filename_only)[0].replace("_", " ").replace("-", " ")
                combined_text = f"[SOURCE FILE: {filename_base}]\n\n{file_text}"
                documents.append(Document(page_content=combined_text, metadata={"source": filepath}))

        except Exception as e:
            print(f"Error loading {filename}: {e}")
    return documents
